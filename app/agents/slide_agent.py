"""
Slide Agent - Creates PowerPoint presentations from research paper content
"""

import logging
from pathlib import Path
from typing import Dict, List

from langchain.schema import HumanMessage, SystemMessage
from langchain_openai import AzureChatOpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


class SlideAgent:
    """Agent responsible for creating PowerPoint slides from research content"""

    def __init__(
        self, azure_endpoint: str, api_key: str, deployment_name: str = "gpt-4o-mini"
    ):
        self.llm = AzureChatOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=api_key,
            azure_deployment=deployment_name,
            openai_api_version="2024-02-01",
            temperature=0.3,
        )

    def generate_slide_content(self, paper_data: Dict) -> List[Dict]:
        """Generate professional slide content structure from paper data"""
        system_prompt = """You are an expert presentation designer and academic content creator.
        Create a structured outline for a 2-3 slide professional presentation based on the research paper.

        Requirements for each slide:
        1. Professional, academic tone with detailed content
        2. Left-aligned text layout for readability
        3. Rich information density - each slide should be substantial
        4. Include specific data, statistics, and technical details
        5. Use bullet points with detailed explanations
        6. Include methodology details, key findings, and implications
        7. Make slides information-rich and visually appealing

        Slide Structure:
        - Slide 1: Research Overview & Methodology (detailed background, objectives, approach)
        - Slide 2: Key Findings & Results (specific data, statistics, significant results)
        - Slide 3: Implications & Conclusions (practical applications, future work, takeaways)

        Each slide should have:
        - A clear, descriptive title
        - 4-6 detailed bullet points with substantial content
        - Specific numbers, percentages, and statistics where available
        - Technical terms with explanations
        - Professional academic language

        Return a JSON list of exactly 2-3 slides with title and bullet_points fields.
        Each bullet point should be detailed and informative (50-100 words each)."""

        sections = paper_data.get("sections", {})
        structure = paper_data.get("structure", {})
        visual_elements = paper_data.get("figures", [])

        # Enhanced content summary with visual elements
        content_summary = f"""
        Paper Analysis: {structure.get("analysis", "")}

        Abstract: {sections.get("abstract", "")[:800]}
        Introduction: {sections.get("introduction", "")[:800]}
        Methods: {sections.get("methods", "")[:800]}
        Results: {sections.get("results", "")[:800]}
        Conclusion: {sections.get("conclusion", "")[:800]}
        Key Findings: {sections.get("key_findings", "")[:600]}
        Statistics: {sections.get("statistics", "")[:600]}
        Methodology: {sections.get("methodology", "")[:600]}
        Implications: {sections.get("implications", "")[:600]}

        Visual Elements Available: {len(visual_elements)} figures/tables
        Large Visual Elements: {len([v for v in visual_elements if v.get("is_large", False)])}
        """

        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=content_summary),
        ]

        try:
            response = self.llm.invoke(messages)
            logger.info("Successfully generated professional slide content structure")
            return self._parse_slide_content(response.content)
        except Exception as e:
            logger.error(f"Error generating slide content: {e}")
            return self._create_professional_default_slides(paper_data)

    def _parse_slide_content(self, content: str) -> List[Dict]:
        """Parse LLM response into slide structure"""
        # Simple parsing - in production, you'd want more robust JSON parsing
        slides = []

        # Default slide structure if parsing fails
        default_slides = [
            {
                "title": "Research Overview",
                "bullet_points": [
                    "Introduction to the research topic",
                    "Key objectives and goals",
                    "Research methodology overview",
                    "Expected outcomes",
                ],
            },
            {
                "title": "Key Findings",
                "bullet_points": [
                    "Main research results",
                    "Statistical significance",
                    "Practical implications",
                    "Future research directions",
                ],
            },
        ]

        try:
            # Try to extract JSON from response
            import json
            import re

            # Look for JSON in the response
            json_match = re.search(r"\[.*\]", content, re.DOTALL)
            if json_match:
                slides = json.loads(json_match.group())
            else:
                slides = default_slides

        except Exception as e:
            logger.warning(f"Could not parse slide content, using defaults: {e}")
            slides = default_slides

        return slides[:3]  # Limit to 3 slides max for POC

    def _create_professional_default_slides(self, paper_data: Dict) -> List[Dict]:
        """Create professional default slides when LLM fails"""
        sections = paper_data.get("sections", {})

        slides = []

        # Slide 1: Research Overview & Methodology
        if sections.get("abstract") or sections.get("introduction"):
            abstract_text = sections.get("abstract", "")[:400]
            intro_text = sections.get("introduction", "")[:400]
            methodology_text = sections.get("methodology", "")[:300]

            slides.append(
                {
                    "title": "Research Overview & Methodology",
                    "bullet_points": [
                        f"Research Focus: {abstract_text.split('.')[0] if abstract_text else 'Academic research study'}",
                        f"Background & Objectives: {intro_text.split('.')[0] if intro_text else 'Comprehensive investigation'}",
                        f"Methodological Approach: {methodology_text.split('.')[0] if methodology_text else 'Systematic research methodology'}",
                        f"Study Design: {sections.get('methods', '')[:200] if sections.get('methods') else 'Rigorous experimental design'}",
                        f"Data Collection: {sections.get('statistics', '')[:200] if sections.get('statistics') else 'Comprehensive data analysis'}",
                    ],
                }
            )

        # Slide 2: Key Findings & Results
        if sections.get("results") or sections.get("key_findings"):
            results_text = sections.get("results", "")[:400]
            findings_text = sections.get("key_findings", "")[:400]
            stats_text = sections.get("statistics", "")[:300]

            slides.append(
                {
                    "title": "Key Findings & Statistical Results",
                    "bullet_points": [
                        f"Primary Results: {results_text.split('.')[0] if results_text else 'Significant research findings'}",
                        f"Statistical Significance: {stats_text.split('.')[0] if stats_text else 'Statistically significant outcomes'}",
                        f"Key Discoveries: {findings_text.split('.')[0] if findings_text else 'Important research discoveries'}",
                        f"Data Analysis: {sections.get('results', '')[:200] if sections.get('results') else 'Comprehensive data interpretation'}",
                        f"Research Impact: {sections.get('implications', '')[:200] if sections.get('implications') else 'Meaningful research contributions'}",
                    ],
                }
            )

        # Slide 3: Implications & Conclusions
        if sections.get("conclusion") or sections.get("implications"):
            conclusion_text = sections.get("conclusion", "")[:400]
            implications_text = sections.get("implications", "")[:400]

            slides.append(
                {
                    "title": "Implications & Future Directions",
                    "bullet_points": [
                        f"Main Conclusions: {conclusion_text.split('.')[0] if conclusion_text else 'Key research conclusions'}",
                        f"Practical Implications: {implications_text.split('.')[0] if implications_text else 'Real-world applications'}",
                        f"Research Contributions: {sections.get('conclusion', '')[:200] if sections.get('conclusion') else 'Significant academic contributions'}",
                        f"Future Research: {sections.get('implications', '')[:200] if sections.get('implications') else 'Directions for future investigation'}",
                        f"Policy Recommendations: {sections.get('implications', '')[:200] if sections.get('implications') else 'Evidence-based recommendations'}",
                    ],
                }
            )

        # If no content available, create basic professional slides
        if not slides:
            slides = [
                {
                    "title": "Research Study Overview",
                    "bullet_points": [
                        "Comprehensive academic research investigation",
                        "Systematic methodology and data collection",
                        "Rigorous analysis and interpretation",
                        "Significant findings and contributions",
                        "Practical implications and applications",
                    ],
                },
                {
                    "title": "Key Research Findings",
                    "bullet_points": [
                        "Statistically significant research outcomes",
                        "Important discoveries and insights",
                        "Data-driven conclusions and interpretations",
                        "Research impact and contributions",
                        "Evidence-based recommendations",
                    ],
                },
                {
                    "title": "Implications & Future Work",
                    "bullet_points": [
                        "Practical applications and implications",
                        "Policy recommendations and guidelines",
                        "Future research directions",
                        "Long-term impact assessment",
                        "Continued investigation opportunities",
                    ],
                },
            ]

        return slides[:3]  # Limit to 3 slides max

    def create_presentation(
        self,
        slides_data: List[Dict],
        output_path: str,
        visual_elements: List[Dict] = None,
    ) -> str:
        """Create professional PowerPoint presentation with left-aligned text and visual elements"""
        try:
            # Create presentation
            prs = Presentation()

            # Set slide size to 16:9
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Get visual elements for inclusion
            if visual_elements is None:
                visual_elements = []

            # Filter for large, important visual elements
            important_visuals = [
                v for v in visual_elements if v.get("is_large", False)
            ][:3]

            for slide_idx, slide_info in enumerate(slides_data):
                # Add slide with title and content layout
                slide_layout = prs.slide_layouts[1]  # Title and content layout
                slide = prs.slides.add_slide(slide_layout)

                # Set title with professional styling
                title = slide.shapes.title
                title.text = slide_info["title"]
                title.text_frame.paragraphs[0].font.size = Pt(36)
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                # Set content with left alignment
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()

                # Add bullet points with left alignment
                for i, bullet_point in enumerate(slide_info["bullet_points"]):
                    p = (
                        text_frame.paragraphs[0]
                        if not text_frame.paragraphs
                        else text_frame.add_paragraph()
                    )
                    p.text = bullet_point
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(51, 51, 51)
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT

                    # Add some spacing between bullet points
                    p.space_after = Pt(12)

                # Add visual element if available for this slide
                if slide_idx < len(important_visuals):
                    visual = important_visuals[slide_idx]
                    self._add_visual_element_to_slide(slide, visual, slide_idx)

            # Save presentation
            prs.save(output_path)
            logger.info(
                f"Created professional presentation with {len(slides_data)} slides: {output_path}"
            )

            return output_path

        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise

    def _add_visual_element_to_slide(self, slide, visual_element: Dict, slide_idx: int):
        """Add visual element (figure/table) to slide"""
        try:
            # Position visual element on the right side of the slide
            left = Inches(7)  # Right side of slide
            top = Inches(1.5)
            width = Inches(5.5)
            height = Inches(4)

            if visual_element.get("type") == "table" and visual_element.get("data"):
                # Add table as text box
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.clear()

                # Add table title
                p = text_frame.paragraphs[0]
                p.text = f"Table {visual_element.get('index', 1)}"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 102)

                # Add table data
                for row in visual_element["data"][:5]:  # Limit to 5 rows
                    p = text_frame.add_paragraph()
                    p.text = row[:50] + "..." if len(row) > 50 else row
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(51, 51, 51)

            elif visual_element.get("data"):
                # Add image placeholder (in a real implementation, you'd embed the actual image)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.clear()

                p = text_frame.paragraphs[0]
                p.text = f"Figure {visual_element.get('index', 1)}"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 102)

                p = text_frame.add_paragraph()
                p.text = visual_element.get(
                    "description", "Visual element from research paper"
                )
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(51, 51, 51)

        except Exception as e:
            logger.warning(f"Could not add visual element to slide: {e}")

    def process_paper_to_slides(
        self, paper_data: Dict, output_dir: str, research_id: int
    ) -> str:
        """Main method to convert paper data to professional PowerPoint slides"""
        logger.info(f"Creating professional slides for research ID: {research_id}")

        # Generate slide content
        slides_data = self.generate_slide_content(paper_data)

        # Create output path
        output_path = Path(output_dir) / f"slides_{research_id}.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Get visual elements for inclusion
        visual_elements = paper_data.get("figures", [])

        # Create presentation with visual elements
        presentation_path = self.create_presentation(
            slides_data, str(output_path), visual_elements
        )

        return presentation_path
