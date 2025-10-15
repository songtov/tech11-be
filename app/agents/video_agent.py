"""
Video Agent - Assembles final video from slides and audio using MoviePy
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional

try:
    from moviepy.editor import (  # type: ignore
        AudioFileClip,
        ImageClip,
        VideoFileClip,
        concatenate_videoclips,
    )
except ImportError:
    # Fallback for moviepy 2.x
    from moviepy import (  # type: ignore
        AudioFileClip,
        ImageClip,
        VideoFileClip,
        concatenate_videoclips,
    )

try:
    from pytoon.animator import animate
except ImportError:
    animate = None

import os
import tempfile

from pptx import Presentation

logger = logging.getLogger(__name__)


class VideoAgent:
    """Agent responsible for creating final video from slides and audio"""

    def __init__(self, fps: int = 1, resolution: tuple = (1920, 1080)):
        self.fps = fps
        self.resolution = resolution

    def convert_pptx_to_images(
        self,
        pptx_path: str,
        output_dir: str,
        generated_figures: List[Optional[str]] = None,
    ) -> List[str]:
        """Convert PowerPoint slides to professional images with generated figures"""
        try:
            # This is a simplified approach - in production you'd use python-pptx with PIL
            # or a more robust solution like LibreOffice headless

            # For now, we'll create placeholder images
            # In a real implementation, you'd extract slides as images
            image_paths = []

            prs = Presentation(pptx_path)

            # Get generated figures for inclusion
            if generated_figures is None:
                generated_figures = []

            for i, slide in enumerate(prs.slides):
                # Create a placeholder image path
                image_path = Path(output_dir) / f"slide_{i + 1:02d}.png"
                image_path.parent.mkdir(parents=True, exist_ok=True)

                # Extract slide content and create image
                slide_content = self._extract_slide_content(slide)

                # Get figure path for this slide
                figure_path = (
                    generated_figures[i] if i < len(generated_figures) else None
                )

                self._create_slide_image(str(image_path), slide_content, figure_path)
                image_paths.append(str(image_path))

            logger.info(f"Converted {len(image_paths)} slides to professional images")
            return image_paths

        except Exception as e:
            logger.error(f"Error converting PPTX to images: {e}")
            raise

    def _extract_slide_content(self, slide):
        """Extract content from a PowerPoint slide"""
        content = {"title": "", "bullet_points": []}

        try:
            # Extract all text from slide shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if not content["title"]:  # First text is usually title
                        content["title"] = text
                    else:
                        # Extract bullet points from content shape
                        lines = text.split("\n")
                        for line in lines:
                            line = line.strip()
                            if line and not line.startswith(
                                "•"
                            ):  # Skip empty lines and bullet markers
                                content["bullet_points"].append(line)
                        break  # Found content, stop looking
        except Exception as e:
            logger.warning(f"Error extracting slide content: {e}")

        return content

    def _create_slide_image(
        self, image_path: str, slide_content: Dict, figure_path: Optional[str] = None
    ):
        """Create a professional slide image with content and generated figure"""
        try:
            from PIL import Image, ImageDraw, ImageFont

            # Create image with professional background
            img = Image.new("RGB", self.resolution, color="white")
            draw = ImageDraw.Draw(img)

            # Load fonts with appropriate sizes
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 60)
                bullet_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 36)
            except (OSError, IOError):
                title_font = ImageFont.load_default()
                bullet_font = ImageFont.load_default()

            # Draw title (left-aligned)
            title = slide_content.get("title", "Untitled Slide")
            title_x = 80  # Left margin
            title_y = 60

            draw.text((title_x, title_y), title, fill=(0, 51, 102), font=title_font)

            # Draw bullet points (left-aligned) with proper spacing
            bullet_points = slide_content.get("bullet_points", [])
            bullet_y = title_y + 100
            max_width = int(self.resolution[0] * 0.55)  # Use 55% of width for text
            line_height = 45  # Line height for text

            for i, point in enumerate(bullet_points[:5]):  # Limit to 5 bullet points
                if (
                    bullet_y > self.resolution[1] - 250
                ):  # Leave space for visual element
                    break

                # Wrap text into lines
                wrapped_lines = self._wrap_text_to_lines(point, bullet_font, max_width)

                # Draw each line of the bullet point
                for line_idx, line in enumerate(wrapped_lines):
                    if bullet_y > self.resolution[1] - 200:
                        break

                    # Add bullet only to first line
                    bullet_text = f"• {line}" if line_idx == 0 else f"  {line}"
                    draw.text(
                        (title_x, bullet_y),
                        bullet_text,
                        fill=(51, 51, 51),
                        font=bullet_font,
                    )
                    bullet_y += line_height

                # Add extra spacing between bullet points
                bullet_y += 15

            # Add generated figure if provided
            if figure_path and Path(figure_path).exists():
                self._add_figure_to_image(img, figure_path)

            # Save image
            img.save(image_path)
            logger.info(f"Created professional slide image: {image_path}")

        except ImportError:
            # Fallback: create a simple text file as placeholder
            with open(image_path.replace(".png", ".txt"), "w") as f:
                f.write(f"Slide: {slide_content.get('title', 'Untitled')}")
        except Exception as e:
            logger.warning(f"Could not create slide image: {e}")

    def _wrap_text_to_lines(self, text: str, font, max_width: int) -> List[str]:
        """Wrap text into lines that fit within specified width"""
        words = text.split()
        lines = []
        current_line = []

        for word in words:
            test_line = " ".join(current_line + [word])
            bbox = font.getbbox(test_line)
            if bbox[2] - bbox[0] <= max_width:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(" ".join(current_line))
                    current_line = [word]
                else:
                    lines.append(word)

        if current_line:
            lines.append(" ".join(current_line))

        return lines[:4]  # Max 4 lines per bullet point

    def _add_figure_to_image(self, img, figure_path: str):
        """Add generated figure to slide image in top-right corner"""
        try:
            from PIL import Image

            # Load the generated figure
            figure_img = Image.open(figure_path)

            # Resize figure to fit in top-right corner (600x400 -> 500x350)
            figure_img = figure_img.resize((500, 350), Image.Resampling.LANCZOS)

            # Position in top-right corner (starting at 60% width)
            right_x = int(self.resolution[0] * 0.6)
            right_y = 80

            # Paste figure onto slide image
            img.paste(figure_img, (right_x, right_y))

            logger.info(f"Added figure to slide at position ({right_x}, {right_y})")

        except Exception as e:
            logger.warning(f"Could not add figure to image: {e}")

    def create_slide_video_clips(
        self, image_paths: List[str], audio_path: str, slide_durations: List[float]
    ) -> List[ImageClip]:
        """Create video clips for each slide"""
        clips = []

        try:
            # Load audio to get total duration
            audio_clip = AudioFileClip(audio_path)
            total_duration = audio_clip.duration

            # Calculate duration per slide
            if slide_durations:
                durations = slide_durations
            else:
                # Equal duration for all slides
                duration_per_slide = total_duration / len(image_paths)
                durations = [duration_per_slide] * len(image_paths)

            for i, image_path in enumerate(image_paths):
                if os.path.exists(image_path):
                    # Create image clip
                    clip = ImageClip(image_path, duration=durations[i])
                    clip = clip.resized(self.resolution)
                    clips.append(clip)
                else:
                    logger.warning(f"Image not found: {image_path}")

            audio_clip.close()

        except Exception as e:
            logger.error(f"Error creating slide video clips: {e}")
            raise

        return clips

    def assemble_video(
        self, video_clips: List[ImageClip], audio_path: str, output_path: str
    ) -> str:
        """Assemble final video from clips and audio"""
        try:
            # Concatenate video clips
            final_video = concatenate_videoclips(video_clips, method="compose")

            # Add audio
            audio_clip = AudioFileClip(audio_path)
            final_video = final_video.with_audio(audio_clip)

            # Write video file
            final_video.write_videofile(
                output_path,
                fps=self.fps,
                codec="libx264",
                audio_codec="aac",
                temp_audiofile="temp-audio.m4a",
                remove_temp=True,
            )

            # Clean up
            final_video.close()
            audio_clip.close()

            logger.info(f"Created final video: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Error assembling video: {e}")
            raise

    def create_video_from_slides_and_audio(
        self,
        pptx_path: str,
        audio_path: str,
        output_path: str,
        slide_durations: Optional[List[float]] = None,
        generated_figures: Optional[List[Optional[str]]] = None,
    ) -> str:
        """Main method to create professional video from slides and audio"""
        logger.info(f"Creating professional video from slides: {pptx_path}")

        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            # Convert slides to images with generated figures
            image_paths = self.convert_pptx_to_images(
                pptx_path, temp_dir, generated_figures
            )

            # Create video clips
            video_clips = self.create_slide_video_clips(
                image_paths, audio_path, slide_durations
            )

            # Assemble final video
            final_video_path = self.assemble_video(video_clips, audio_path, output_path)

            return final_video_path

    def process_slides_and_audio_to_video(
        self,
        slides_path: str,
        audio_data: Dict,
        output_dir: str,
        research_id: int,
        generated_figures: List[Optional[str]] = None,
    ) -> str:
        """Main method to process slides and audio into final professional video with animation overlay"""
        logger.info(
            f"Creating final professional video with animation for research ID: {research_id}"
        )

        # Create output paths
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        # Base video path (without animation)
        base_video_path = output_path / f"base_video_{research_id}.mp4"
        # Final video path (with animation overlay)
        final_video_path = output_path / f"video_{research_id}.mp4"

        # Get audio and transcript paths
        audio_path = audio_data["full_audio_path"]
        transcript_path = audio_data.get("transcript_path")

        # Calculate slide durations based on audio segments
        slide_durations = []
        if "slide_audio_files" in audio_data:
            for audio_file in audio_data["slide_audio_files"]:
                try:
                    audio_clip = AudioFileClip(audio_file)
                    duration = audio_clip.duration
                    audio_clip.close()
                    slide_durations.append(duration)
                except Exception as e:
                    logger.warning(f"Could not get duration for {audio_file}: {e}")
                    slide_durations.append(5.0)  # Default duration

        # Step 1: Create base slide video (existing functionality)
        logger.info("Step 1: Creating base slide video...")
        base_video_path_str = self.create_video_from_slides_and_audio(
            slides_path,
            audio_path,
            str(base_video_path),
            slide_durations,
            generated_figures,
        )

        # Step 2: Add PyToon animation overlay if transcript is available
        if transcript_path and Path(transcript_path).exists():
            logger.info("Step 2: Adding PyToon animation overlay...")
            final_video_path_str = self._add_animation_overlay(
                str(base_video_path), audio_path, transcript_path, str(final_video_path)
            )
        else:
            logger.warning(
                "Transcript not available, using base video without animation"
            )
            # Copy base video to final path if no animation
            import shutil
            shutil.copy2(base_video_path_str, str(final_video_path))
            final_video_path_str = str(final_video_path)

        # Clean up temporary base video file
        try:
            if (
                base_video_path.exists()
                and str(base_video_path) != final_video_path_str
            ):
                base_video_path.unlink()
                logger.info("Cleaned up temporary base video file")
        except Exception as e:
            logger.warning(f"Could not clean up base video file: {e}")

        return final_video_path_str

    def _add_animation_overlay(
        self,
        base_video_path: str,
        audio_path: str,
        transcript_path: str,
        output_path: str,
    ) -> str:
        """Add PyToon animated character overlay to the base slide video"""
        try:
            if animate is None:
                logger.warning(
                    "PyToon not available, copying base video without animation"
                )
                import shutil
                shutil.copy2(base_video_path, output_path)
                return output_path

            logger.info("Creating PyToon animation overlay...")

            # Read transcript content
            with open(transcript_path, "r", encoding="utf-8") as file:
                transcript = file.read()

            # Create PyToon animation
            logger.info("Generating animated character with lip-sync...")
            animation = animate(
                audio_file=audio_path,
                transcript=transcript,
            )

            # Load base video as background
            background_video = VideoFileClip(base_video_path)

            # Export animation with base video as background, positioned bottom-down
            logger.info("Overlaying animation on base video...")
            animation.export(path=output_path, background=background_video, scale=0.7)

            # Clean up
            background_video.close()

            logger.info(
                f"Successfully created video with animated character: {output_path}"
            )
            return output_path

        except Exception as e:
            logger.error(f"Error adding animation overlay: {e}")
            logger.warning("Falling back to base video without animation")
            import shutil
            shutil.copy2(base_video_path, output_path)
            return output_path
